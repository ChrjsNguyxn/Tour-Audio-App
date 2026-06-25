# -*- coding: utf-8 -*-
"""
Kinh tế Chính trị Mác-Lênin: Cách mạng Công nghiệp & Công nghiệp hóa
Tông màu: Vàng nâu nhạt học thuật (Beige & Sand Scholarly) từ ảnh mẫu
Yêu cầu: Thư viện python-pptx (Cài đặt bằng lệnh: pip install python-pptx)
"""

import os
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Vui lòng cài đặt thư viện python-pptx bằng lệnh: pip install python-pptx")
    raise

# 1. HỆ MÀU VÀNG NÂU HỌC THUẬT (Đồng bộ chuẩn xác với ảnh mẫu)
BG_CREAM = RGBColor(0xF7, 0xF4, 0xEE)     # Nền chính kem ngà nhã nhặn
BG_SAND = RGBColor(0xEA, 0xE5, 0xDA)      # Vàng cát ấm cho các trang chuyển phần
TEXT_DARK = RGBColor(0x2D, 0x29, 0x26)    # Chữ xám đen cổ điển ấm áp
ACCENT_GOLD = RGBColor(0x8C, 0x62, 0x39)  # Vàng đồng / Nâu đồng cổ kính làm điểm nhấn
BORDER_COLOR = RGBColor(0xD2, 0xC4, 0xB1) # Đường viền vàng cát sáng thanh lịch
WHITE = RGBColor(0xFF, 0xFF, 0xFF)        # Nền trắng cho các thẻ nội dung
ACCENT_RED = RGBColor(0xA8, 0x44, 0x32)   # Đỏ gạch trầm nhã nhặn cho Thách thức
ACCENT_GREEN = RGBColor(0x4E, 0x7A, 0x5A) # Xanh rêu nhẹ nhàng cho Xu hướng

# Khởi tạo bản trình chiếu
prs = Presentation()
prs.slide_width = Inches(13.333) # Tỷ lệ chuẩn 16:9
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

# --- CÁC HÀM TRỢ GIÚP THIẾT KẾ ---

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_header(slide, title_text):
    # Tiêu đề slide chuẩn góc trên lề trái
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Georgia'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_DARK
    
    # Đường kẻ ngang màu nâu đồng mỏng sang trọng
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(11.733), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_GOLD
    line.line.color.rgb = ACCENT_GOLD

def add_card(slide, left, top, width, height, title, content, border_color=BORDER_COLOR):
    # Tạo hộp nội dung bo góc thanh nhã
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = border_color
    card.line.width = Pt(1.5)
    
    # Tiêu đề hộp
    tb_title = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.5))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.name = 'Georgia'
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = ACCENT_GOLD
    
    # Nội dung hộp
    tb_content = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.75), width - Inches(0.4), height - Inches(0.9))
    tf_content = tb_content.text_frame
    tf_content.word_wrap = True
    p_content = tf_content.paragraphs[0]
    p_content.text = content
    p_content.font.name = 'Calibri'
    p_content.font.size = Pt(14)
    p_content.font.color.rgb = TEXT_DARK

def add_image_placeholder(slide, left, top, width, height, desc_text):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = BG_SAND
    box.line.color.rgb = BORDER_COLOR
    box.line.width = Pt(1)
    
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = f"[ Khung ảnh minh họa ]\n\n({desc_text})"
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.italic = True
    p.font.color.rgb = ACCENT_GOLD


# ==============================================================================
# SLIDE 1: TRANG TIÊU ĐỀ
# ==============================================================================
s1 = prs.slides.add_slide(blank_layout)
set_slide_background(s1, BG_CREAM)

# Đường gờ trang trí bên trái lề giống slide mẫu
girdle = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
girdle.fill.solid()
girdle.fill.fore_color.rgb = ACCENT_GOLD
girdle.line.fill.background()

title_box = s1.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.33), Inches(4.5))
tf = title_box.text_frame
tf.word_wrap = True

p0 = tf.paragraphs[0]
p0.text = "HỌC PHẦN: KINH TẾ CHÍNH TRỊ MÁC - LÊNIN"
p0.font.name = 'Calibri'
p0.font.size = Pt(16)
p0.font.bold = True
p0.font.color.rgb = ACCENT_GOLD
p0.space_after = Pt(20)

p1 = tf.add_paragraph()
p1.text = "CÁCH MẠNG CÔNG NGHIỆP\nVÀ CÔNG NGHIỆP HÓA"
p1.font.name = 'Georgia'
p1.font.size = Pt(44)
p1.font.bold = True
p1.font.color.rgb = TEXT_DARK
p1.space_after = Pt(30)

p2 = tf.add_paragraph()
p2.text = "Nhóm thuyết trình: Nguyễn Quốc Minh  |  Phạm vi: Giáo trình Trang 225 – 245"
p2.font.name = 'Calibri'
p2.font.size = Pt(16)
p2.font.italic = True
p2.font.color.rgb = TEXT_DARK


# ==============================================================================
# SLIDE 2: SECTION TITLE (Phần 1: Tóm tắt ý chính)
# ==============================================================================
s2 = prs.slides.add_slide(blank_layout)
set_slide_background(s2, BG_SAND)

divider_box = s2.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.33), Inches(3.0))
tf2 = divider_box.text_frame
tf2.word_wrap = True

p_div_sub = tf2.paragraphs[0]
p_div_sub.alignment = PP_ALIGN.CENTER
p_div_sub.text = "PHẦN 1"
p_div_sub.font.name = 'Calibri'
p_div_sub.font.size = Pt(18)
p_div_sub.font.bold = True
p_div_sub.font.color.rgb = ACCENT_GOLD
p_div_sub.space_after = Pt(10)

p_div_main = tf2.add_paragraph()
p_div_main.alignment = PP_ALIGN.CENTER
p_div_main.text = "TÓM TẮT Ý CHÍNH CỐT LÕI"
p_div_main.font.name = 'Georgia'
p_div_main.font.size = Pt(40)
p_div_main.font.bold = True
p_div_main.font.color.rgb = TEXT_DARK
p_div_main.space_after = Pt(15)

p_div_desc = tf2.add_paragraph()
p_div_desc.alignment = PP_ALIGN.CENTER
p_div_desc.text = "Khái quát các khái niệm, tiến trình lịch sử và thực tiễn Công nghiệp hóa tại Việt Nam"
p_div_desc.font.name = 'Calibri'
p_div_desc.font.size = Pt(16)
p_div_desc.font.color.rgb = TEXT_DARK


# ==============================================================================
# SLIDE 3: CÁC KHÁI NIỆM CƠ BẢN (3 Cards)
# ==============================================================================
s3 = prs.slides.add_slide(blank_layout)
set_slide_background(s3, BG_CREAM)
add_slide_header(s3, "Các khái niệm nền tảng trong học phần")

add_card(s3, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8), 
         "1. CÔNG NGHIỆP", 
         "Là ngành sản xuất vật chất thuộc khu vực II, sử dụng máy móc thiết bị và công nghệ hiện đại để khai thác, chế biến tài nguyên thiên nhiên tạo ra của cải vật chất phục vụ đời sống xã hội.")

add_card(s3, Inches(4.86), Inches(1.8), Inches(3.6), Inches(4.8), 
         "2. CÔNG NGHIỆP HÓA", 
         "Là quá trình chuyển đổi căn bản, toàn diện các hoạt động sản xuất, kinh doanh, dịch vụ và quản lý từ sử dụng lao động thủ công là chính sang lao động sử dụng máy móc và công nghệ tiên tiến nhằm tạo ra năng suất lao động xã hội cao vượt trội.")

add_card(s3, Inches(8.92), Inches(1.8), Inches(3.6), Inches(4.8), 
         "3. CÁCH MẠNG CÔNG NGHIỆP", 
         "Là bước nhảy vọt vĩ đại về chất trong lực lượng sản xuất và tư liệu lao động, dựa trên những phát minh mang tính đột phá về khoa học và kỹ thuật, từ đó thay đổi sâu sắc phương thức sản xuất và mọi mặt đời sống xã hội.")


# ==============================================================================
# SLIDE 4: TIMELINE 4 CUỘC CMCN
# ==============================================================================
s4 = prs.slides.add_slide(blank_layout)
set_slide_background(s4, BG_CREAM)
add_slide_header(s4, "Tiến trình các cuộc Cách mạng Công nghiệp")

line_timeline = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.85), Inches(11.733), Inches(0.04))
line_timeline.fill.solid()
line_timeline.fill.fore_color.rgb = BORDER_COLOR
line_timeline.line.fill.background()

col_w = Inches(2.65)
gap = Inches(0.35)

milestones = [
    ("CMCN 1.0 (Thế kỷ XVIII)", "Cơ giới hóa\nKhởi nguồn từ dệt may nước Anh. Máy hơi nước của James Watt thay thế sức người bằng cơ khí."),
    ("CMCN 2.0 (Thế kỷ XIX)", "Điện khí hóa\nSự ra đời của điện năng, động cơ đốt trong và dây chuyền lắp ráp hàng loạt nâng cao năng suất."),
    ("CMCN 3.0 (Thập niên 1960)", "Tự động hóa\nSự bùng nổ của chất bán dẫn, máy tính cá nhân và Internet giúp tự động hóa dây chuyền."),
    ("CMCN 4.0 (Hiện nay)", "Thông minh hóa\nSự kết hợp giữa AI, IoT, Big Data tạo nên hệ thống sản xuất siêu kết nối thông minh.")
]

for i, (title, text) in enumerate(milestones):
    col_left = Inches(0.8) + i * (col_w + gap)
    card = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_left, Inches(1.8) if i%2==0 else Inches(4.2), col_w, Inches(2.2))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = ACCENT_GOLD if i == 3 else BORDER_COLOR
    card.line.width = Pt(1.5)
    
    tb = s4.shapes.add_textbox(col_left + Inches(0.15), (Inches(1.9) if i%2==0 else Inches(4.3)), col_w - Inches(0.3), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.name = 'Georgia'
    p1.font.size = Pt(15)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_GOLD
    p1.space_after = Pt(4)
    
    p2 = tf.add_paragraph()
    p2.text = text
    p2.font.name = 'Calibri'
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_DARK


# ==============================================================================
# SLIDE 5: CÔNG NGHIỆP HÓA Ở VIỆT NAM (2 Cột)
# ==============================================================================
s5 = prs.slides.add_slide(blank_layout)
set_slide_background(s5, BG_CREAM)
add_slide_header(s5, "Tính tất yếu khách quan của CNH ở Việt Nam")

left_box = s5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.0), Inches(5.0))
tf_left = left_box.text_frame
tf_left.word_wrap = True

p_main = tf_left.paragraphs[0]
p_main.text = "Tại sao Việt Nam bắt buộc phải tiến hành CNH?"
p_main.font.name = 'Georgia'
p_main.font.size = Pt(22)
p_main.font.bold = True
p_main.font.color.rgb = TEXT_DARK
p_main.space_after = Pt(15)

bullets = [
    "Xây dựng cơ sở vật chất - kỹ thuật cho CNXH: Việt Nam quá độ từ nền nông nghiệp lạc hậu, bắt buộc phải CNH để hiện đại hóa sản xuất xã hội.",
    "Phát triển lực lượng sản xuất: Nâng cao năng suất lao động tổng thể, tăng khả năng cạnh tranh trong kỷ nguyên kinh tế số.",
    "Bảo vệ Tổ quốc vững chắc: Tạo cơ sở vật chất củng cố tiềm lực quốc phòng toàn dân và an ninh quốc gia chủ động.",
    "Nâng cao vị thế kinh tế: Hội nhập kinh tế quốc tế sâu rộng mà vẫn giữ vững tính độc lập, tự chủ."
]

for bullet in bullets:
    p = tf_left.add_paragraph()
    p.text = "• " + bullet
    p.font.name = 'Calibri'
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_DARK
    p.space_after = Pt(10)

add_image_placeholder(s5, Inches(7.4), Inches(1.8), Inches(5.1), Inches(4.8), 
                      "Có thể chèn hình ảnh Nhà máy dệt thời xưa hoặc Cơ khí hóa Việt Nam")


# ==============================================================================
# SLIDE 6: 3 NỘI DUNG CỐT LÕI
# ==============================================================================
s6 = prs.slides.add_slide(blank_layout)
set_slide_background(s6, BG_CREAM)
add_slide_header(s6, "3 Nội dung cốt lõi của CNH ở Việt Nam")

add_card(s6, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8),
         "I. CHUYỂN DỊCH CƠ CẤU KINH TẾ",
         "• Giảm tỷ trọng nông nghiệp nông thôn; tăng nhanh tỷ trọng công nghiệp chế tạo và dịch vụ hiện đại.\n• Quy hoạch phát triển vùng kinh tế trọng điểm dựa theo lợi thế cạnh tranh tự nhiên của từng vùng miền.")

add_card(s6, Inches(4.86), Inches(1.8), Inches(3.6), Inches(4.8),
         "II. PHÁT TRIỂN LỰC LƯỢNG SẢN XUẤT",
         "• Chuyển giao, ứng dụng mạnh mẽ các đột phá công nghệ mới vào sản xuất.\n• Coi trọng giáo dục đào tạo nguồn nhân lực chất lượng cao – chìa khóa quyết định sự thịnh vượng bền vững.")

add_card(s6, Inches(8.92), Inches(1.8), Inches(3.6), Inches(4.8),
         "III. HOÀN THIỆN QUAN HỆ SẢN XUẤT",
         "• Xây dựng đồng bộ thể chế kinh tế thị trường định hướng xã hội chủ nghĩa.\n• Phát triển kinh tế nhiều thành phần, giữ vững vai trò dẫn dắt chủ đạo của thành phần kinh tế nhà nước.")


# ==============================================================================
# SLIDE 7: THÁCH THỨC VÀ XU HƯỚNG HIỆN ĐẠI
# ==============================================================================
s7 = prs.slides.add_slide(blank_layout)
set_slide_background(s7, BG_CREAM)
add_slide_header(s7, "Thách thức & Xu hướng phát triển hiện đại")

add_card(s7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), 
         "⚠ THÁCH THỨC LỚN NHẤT", 
         "1. Già hóa dân số nhanh:\nDự kiến đạt tỷ lệ 25% người cao tuổi vào năm 2050, gây áp lực trực tiếp lên lực lượng lao động trẻ.\n\n2. Bẫy thu nhập trung bình:\nNguy cơ tụt hậu sâu sắc về năng suất lao động xã hội và công nghệ lõi tự chủ.\n\n3. Chảy máu chất xám:\nThách thức thu hút và đãi ngộ nhân tài kỹ thuật trình độ cao.", 
         border_color=ACCENT_RED)

add_card(s7, Inches(6.933), Inches(1.8), Inches(5.6), Inches(4.8), 
         "☘ XU HƯỚNG CHUYỂN DỊCH", 
         "1. Chuyển đổi số toàn diện:\nĐưa giải pháp thông minh vào chuỗi vận hành quản trị và dịch vụ công.\n\n2. Kinh tế xanh & tuần hoàn:\nChuyển đổi từ mô hình tuyến tính sang tuần hoàn, cam kết hướng tới Net Zero 2050.\n\n3. Công nghiệp nền tảng:\nTập trung phát triển công nghiệp chế biến, chế tạo nội địa có giá trị gia tăng lớn.", 
         border_color=ACCENT_GREEN)


# ==============================================================================
# SLIDE 8: SECTION TITLE (Phần 2: 10 Câu hỏi trắc nghiệm)
# ==============================================================================
s8 = prs.slides.add_slide(blank_layout)
set_slide_background(s8, BG_SAND)

divider_box = s8.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.33), Inches(3.0))
tf8 = divider_box.text_frame
tf8.word_wrap = True

p_div_sub2 = tf8.paragraphs[0]
p_div_sub2.alignment = PP_ALIGN.CENTER
p_div_sub2.text = "PHẦN 2"
p_div_sub2.font.name = 'Calibri'
p_div_sub2.font.size = Pt(18)
p_div_sub2.font.bold = True
p_div_sub2.font.color.rgb = ACCENT_GOLD
p_div_sub2.space_after = Pt(10)

p_div_main2 = tf8.add_paragraph()
p_div_main2.alignment = PP_ALIGN.CENTER
p_div_main2.text = "10 CÂU HỎI TRẮC NGHIỆM ÔN TẬP"
p_div_main2.font.name = 'Georgia'
p_div_main2.font.size = Pt(40)
p_div_main2.font.bold = True
p_div_main2.font.color.rgb = TEXT_DARK
p_div_main2.space_after = Pt(15)

p_div_desc2 = tf8.add_paragraph()
p_div_desc2.alignment = PP_ALIGN.CENTER
p_div_desc2.text = "Hệ thống câu hỏi ôn tập bám sát nội dung giáo trình chính quy\n(Hướng dẫn cài đặt hiệu ứng click chuột hiện đáp án nằm ở phần Ghi chú dưới Slide)"
p_div_desc2.font.name = 'Calibri'
p_div_desc2.font.size = Pt(14)
p_div_desc2.font.italic = True
p_div_desc2.font.color.rgb = TEXT_DARK


# ==============================================================================
# HÀM TỰ ĐỘNG DÀNG TRANG TRẮC NGHIỆM
# ==============================================================================
def add_mcq_card(slide, left, top, width, height, q_num, question, options, correct_index):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = BORDER_COLOR
    box.line.width = Pt(1.5)
    
    tf = box.text_frame
    tf.word_wrap = True
    
    p_q = tf.paragraphs[0]
    p_q.text = f"Câu {q_num}: {question}"
    p_q.font.name = 'Georgia'
    p_q.font.size = Pt(14)
    p_q.font.bold = True
    p_q.font.color.rgb = TEXT_DARK
    p_q.space_after = Pt(10)
    
    for idx, option in enumerate(options):
        p_opt = tf.add_paragraph()
        p_opt.text = option
        p_opt.font.name = 'Calibri'
        p_opt.font.size = Pt(12)
        if idx == correct_index:
            p_opt.font.bold = True
            p_opt.font.color.rgb = ACCENT_GOLD
            p_opt.text += "  ✓ (Đáp án đúng)"
        else:
            p_opt.font.color.rgb = TEXT_DARK
        p_opt.space_after = Pt(4)

MCQ_NOTES = (
    "HƯỚNG DẪN CÀI ĐẶT HIỆU ỨNG CLICK HIỆN ĐÁP ÁN TRONG POWERPOINT:\n"
    "1. Do tệp tạo bằng Python là tệp tĩnh, bạn hãy nhấp chọn khối câu hỏi này trên trang.\n"
    "2. Trong PowerPoint, bạn hãy vẽ một nút bấm nhỏ đè lên phần chữ 'Đáp án đúng' hoặc tạo một nút 'Hiện đáp án'.\n"
    "3. Chọn dòng chữ đáp án đúng có đánh dấu ✓, chuyển sang tab 'Animations' (Hoạt họa) -> chọn hiệu ứng 'Appear' (Xuất hiện) hoặc 'Fade' (Mờ dần).\n"
    "4. Ở góc phải Tab Animations, chọn 'Trigger' -> 'On Click of' -> chọn đúng tên nút bấm của câu hỏi đó.\n"
    "Chúc bạn có một buổi thuyết trình tương tác thành công rực rỡ!"
)


# ==============================================================================
# SLIDE 9: CÂU HỎI TRẮC NGHIỆM (Câu 1 - 3)
# ==============================================================================
s9 = prs.slides.add_slide(blank_layout)
set_slide_background(s9, BG_CREAM)
add_slide_header(s9, "Câu hỏi ôn tập kiến thức (Câu 1 - 3)")
s9.notes_slide.notes_text_frame.text = MCQ_NOTES

add_mcq_card(s9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.2), 1, 
             "Công nghiệp hóa (CNH) được định nghĩa chính xác nhất là gì?",
             ["A. Tập trung xây dựng nhiều nhà máy xí nghiệp quy mô lớn trên cả nước",
              "B. Chuyển đổi căn bản toàn diện nền sản xuất từ lao động thủ công sang lao động bằng máy móc"], 1)

add_mcq_card(s9, Inches(6.933), Inches(1.8), Inches(5.6), Inches(2.2), 2, 
             "Cuộc Cách mạng Công nghiệp lần thứ nhất bắt đầu từ quốc gia nào?",
             ["A. Nước Anh, khởi đầu từ giữa thế kỷ XVIII",
              "B. Nước Pháp, khởi đầu từ đầu thế kỷ XIX",
              "C. Nước Đức, khởi đầu từ cuối thế kỷ XVII"], 0)

add_mcq_card(s9, Inches(0.8), Inches(4.4), Inches(11.733), Inches(2.3), 3, 
             "Phát minh then chốt nào được xem là mốc mở đầu cho quá trình cơ giới hóa sản xuất?",
             ["A. Máy thoi bay trong ngành dệt của John Kay phát minh vào năm 1733",
              "B. Động cơ hơi nước cải tiến hoàn thiện của James Watt chế tạo vào năm 1784",
              "C. Đầu máy xe lửa chạy bằng hơi nước của Stephenson chế tạo vào năm 1814"], 1)


# ==============================================================================
# SLIDE 10: CÂU HỎI TRẮC NGHIỆM (Câu 4 - 6)
# ==============================================================================
s10 = prs.slides.add_slide(blank_layout)
set_slide_background(s10, BG_CREAM)
add_slide_header(s10, "Câu hỏi ôn tập kiến thức (Câu 4 - 6)")
s10.notes_slide.notes_text_frame.text = MCQ_NOTES

add_mcq_card(s10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(2.2), 4, 
             "Theo học thuyết của C. Mác, cuộc CMCN 1.0 trải qua giai đoạn phát triển cao nhất nào?",
             ["A. Hiệp tác giản đơn thủ công nghiệp",
              "B. Công trường thủ công bán cơ khí",
              "C. Đại công nghiệp cơ khí hóa hoàn chỉnh"], 2)

add_mcq_card(s10, Inches(6.933), Inches(1.8), Inches(5.6), Inches(2.2), 5, 
             "Đặc trưng kỹ thuật nổi bật nhất của Cách mạng Công nghiệp 2.0 là gì?",
             ["A. Động cơ hơi nước và dệt máy cơ khí",
              "B. Điện năng, dây chuyền lắp ráp hàng loạt và động cơ đốt trong",
              "C. Máy vi tính, chất bán dẫn phổ cập và mạng lưới internet"], 1)

add_mcq_card(s10, Inches(0.8), Inches(4.4), Inches(11.733), Inches(2.3), 6, 
             "Điểm cốt lõi nhất giúp phân biệt rõ rệt giữa CMCN 3.0 và CMCN 4.0 là gì?",
             ["A. CMCN 3.0 dựa trên số hóa tự động hóa riêng lẻ; CMCN 4.0 là hệ thống thông minh, siêu kết nối AI, IoT và Big Data thực thể",
              "B. CMCN 3.0 phát triển năng lượng hạt nhân còn CMCN 4.0 tập trung năng lượng hóa thạch",
              "C. Không có sự cải tiến khoa học khác biệt, chỉ là tên gọi thương mại mới thay thế"], 0)


# ==============================================================================
# SLIDE 11: CÂU HỎI TRẮC NGHIỆM (Câu 7 - 10)
# ==============================================================================
s11 = prs.slides.add_slide(blank_layout)
set_slide_background(s11, BG_CREAM)
add_slide_header(s11, "Câu hỏi ôn tập kiến thức (Câu 7 - 10)")
s11.notes_slide.notes_text_frame.text = MCQ_NOTES

col_w2 = Inches(5.6)
row_h2 = Inches(2.2)

add_mcq_card(s11, Inches(0.8), Inches(1.8), col_w2, row_h2, 7, 
             "Đặc điểm nào dưới đây KHÔNG phải là đặc điểm của CNH ở nước ta?",
             ["A. Tiến hành CNH đi kèm mục tiêu định hướng xã hội chủ nghĩa",
              "B. CNH phát triển theo mô hình tự do hoang dã, không có sự định hướng quản lý từ nhà nước"], 1)

add_mcq_card(s11, Inches(6.933), Inches(1.8), col_w2, row_h2, 8, 
             "Thách thức nhân khẩu học nghiêm trọng nhất của CNH Việt Nam là gì?",
             ["A. Tỷ lệ già hóa dân số tăng quá nhanh và nguy cơ thiếu lao động kỹ thuật cao",
              "B. Quy mô dân số giảm sút tuyệt đối đe dọa năng lực sản xuất"], 0)

add_mcq_card(s11, Inches(0.8), Inches(4.4), col_w2, row_h2, 9, 
             "Vì sao Việt Nam bắt buộc phải coi CNH là nhiệm vụ trung tâm suốt thời kỳ quá độ?",
             ["A. Do xuất phát điểm từ nước nông nghiệp lạc hậu, cần xây dựng cơ sở vật chất - kỹ thuật cho CNXH",
              "B. Nhằm phục vụ hoàn hảo các yêu cầu viện trợ của đối tác phương Tây"], 0)

add_mcq_card(s11, Inches(6.933), Inches(4.4), col_w2, row_h2, 10, 
             "Vị trí và vai trò của kinh tế nhà nước trong quá trình CNH là gì?",
             ["A. Giữ vai trò chủ đạo định hướng, hỗ trợ các thành phần kinh tế phát triển",
              "B. Là thành phần kinh tế duy nhất được phép hoạt động"], 0)


# ==============================================================================
# SLIDE 12: KẾT LUẬN / Q&A
# ==============================================================================
s12 = prs.slides.add_slide(blank_layout)
set_slide_background(s12, BG_CREAM)

# Đường gờ lề trái trang nhã
girdle12 = s12.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
girdle12.fill.solid()
girdle12.fill.fore_color.rgb = ACCENT_GOLD
girdle12.line.fill.background()

thanks_box = s12.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10.33), Inches(4.5))
tf12 = thanks_box.text_frame
tf12.word_wrap = True

p_thanks = tf12.paragraphs[0]
p_thanks.alignment = PP_ALIGN.CENTER
p_thanks.text = "XIN CHÂN THÀNH CẢM ƠN!"
p_thanks.font.name = 'Georgia'
p_thanks.font.size = Pt(44)
p_thanks.font.bold = True
p_thanks.font.color.rgb = ACCENT_GOLD
p_thanks.space_after = Pt(15)

p_qa = tf12.add_paragraph()
p_qa.alignment = PP_ALIGN.CENTER
p_qa.text = "Rất mong nhận được sự đóng góp và thảo luận từ Thầy cùng các bạn học viên."
p_qa.font.name = 'Calibri'
p_qa.font.size = Pt(18)
p_qa.font.italic = True
p_qa.font.color.rgb = TEXT_DARK
p_qa.space_after = Pt(40)

contact_card = s12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.66), Inches(4.4), Inches(6.0), Inches(1.5))
contact_card.fill.solid()
contact_card.fill.fore_color.rgb = WHITE
contact_card.line.color.rgb = BORDER_COLOR
contact_card.line.width = Pt(1.5)

tf_c = contact_card.text_frame
tf_c.word_wrap = True
tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = Inches(0.15)

pc1 = tf_c.paragraphs[0]
pc1.alignment = PP_ALIGN.CENTER
pc1.text = "THÔNG TIN LIÊN HỆ & TÀI LIỆU KHẢO SÁT"
pc1.font.name = 'Georgia'
pc1.font.size = Pt(13)
pc1.font.bold = True
pc1.font.color.rgb = ACCENT_GOLD
pc1.space_after = Pt(8)

pc2 = tf_c.add_paragraph()
pc2.alignment = PP_ALIGN.CENTER
pc2.text = "Người thực hiện: Nguyễn Quốc Minh\nNguồn tham chiếu chính quy: Giáo trình Kinh tế Chính trị Mác-Lênin"
pc2.font.name = 'Calibri'
pc2.font.size = Pt(11)
pc2.font.color.rgb = TEXT_DARK

# LƯU CHÍNH XÁC VỚI TÊN FILE MONG MUỐN
output_filename = "CMCN_TomTat_TracNghiem_1.pptx"
try:
    prs.save(output_filename)
    print(f"\n[Thành công] Tệp slide PowerPoint chuyên nghiệp đã được tạo thành công: '{output_filename}'")
except Exception as e:
    print(f"\n[Lỗi] Không thể lưu file: {e}")